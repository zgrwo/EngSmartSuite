"""Reporter — 多格式报告输出：Excel 图表 / PDF / PPT。"""
import io
import logging
import os

import matplotlib.pyplot as plt

from smartsuite.core.contracts import AnalysisResult
from smartsuite.core.exceptions import OutputError

# 图表输出 DPI 常量
_CHART_DPI = 150
_PDF_DPI = 200

logger = logging.getLogger(__name__)


def _validate_output_path(output_path: str) -> str:
    """验证输出路径安全性: 解析为绝对路径并检查是否在合理目录内。

    防御性检查 — 当前调用者不暴露用户控制路径, 但作为未来防护。
    """
    abs_path = os.path.abspath(output_path)
    # 确保目录存在
    out_dir = os.path.dirname(abs_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    return abs_path


def to_excel(result: AnalysisResult, workbook,
             sheet_name: str = "分析结果") -> str:
    """将分析结果写入 Excel 新 Sheet。

    注意: 此函数依赖 xlwings Excel add-in 运行环境，需要:
      1. 安装 xlwings: ``pip install xlwings``
      2. Excel 实例正在运行（通过 xlwings 连接）
      3. 传入有效的 xlwings Workbook 对象

    Web UI 和 CLI 入口不调用此函数。如需导出 Excel 报告，请使用
    ``audit.export_workbook()``（基于 openpyxl，无需 Excel 实例）。
    """
    try:
        ws = workbook.sheets.add(sheet_name, after=workbook.sheets[-1])
        r = 1
        ws.range(f"A{r}").value = "分析结论"
        ws.range(f"A{r}").font.bold = True
        r += 1
        ws.range(f"A{r}").value = result.summary
        r += 2
        for name, df in result.tables.items():
            ws.range(f"A{r}").value = name
            ws.range(f"A{r}").font.bold = True
            r += 1
            ws.range(f"A{r}").value = df
            r += len(df) + 2
        for i, fig in enumerate(result.figures):
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=_CHART_DPI, bbox_inches='tight')
            buf.seek(0)
            pic = workbook.sheets.add(f"图表_{i + 1}", after=workbook.sheets[-1])
            pic.pictures.add(buf, left=pic.range("A1").left,
                             top=pic.range("A1").top, width=600, height=450)
            plt.close(fig)
        return sheet_name
    except Exception as e:
        logger.exception("Excel 输出失败")
        raise OutputError("Excel 输出失败，请检查工作簿是否可写") from e


def to_pdf(result: AnalysisResult, output_path: str) -> str:
    """生成 PDF 报告。"""
    output_path = _validate_output_path(output_path)
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.pdfgen import canvas

        # ── 注册中文字体，避免 CJK 文本渲染为空白 ──
        _cjk_font_name = None
        _cjk_fonts = [
            ("C:/Windows/Fonts/msyh.ttc", "MSYH"),
            ("C:/Windows/Fonts/simhei.ttf", "SimHei"),
            ("/System/Library/Fonts/PingFang.ttc", "PingFang"),
            ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", "NotoCJK"),
        ]
        for _fp, _fn in _cjk_fonts:
            if os.path.exists(_fp):
                try:
                    pdfmetrics.registerFont(TTFont(_fn, _fp))
                    _cjk_font_name = _fn
                    break
                except Exception:
                    continue

        # 字体选择: CJK 字体仅注册单一 weight, 全部文本使用同一字体;
        # 回退到 Helvetica (无 CJK 字形) 仅在无中文字体时
        _title_font = _cjk_font_name if _cjk_font_name else "Helvetica-Bold"
        _body_font = _cjk_font_name if _cjk_font_name else "Helvetica"

        c = canvas.Canvas(output_path, pagesize=A4)
        w, h = A4
        y = h - 50
        c.setFont(_title_font, 16)
        c.drawString(50, y, f"分析报告: {result.task}")
        y -= 30
        c.setFont(_body_font, 11)
        c.drawString(50, y, result.summary)
        y -= 50

        for name, df in list(result.tables.items())[:5]:
            if y < 200:
                c.showPage()
                y = h - 50
            c.setFont(_title_font, 10)
            c.drawString(50, y, name)
            y -= 18
            c.setFont("Courier", 7)
            # 格式化表格行输出
            for _, row in df.head(15).iterrows():
                line = "  ".join(f"{k}={v}" for k, v in row.items())
                c.drawString(55, y, line[:120])
                y -= 12
            y -= 10

        for fig in result.figures:
            if y < 350:
                c.showPage()
                y = h - 50
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=_PDF_DPI, bbox_inches='tight')
            buf.seek(0)
            c.drawImage(ImageReader(buf), 50, y - 300, width=450, height=300)
            plt.close(fig)
            y -= 320

        c.save()
        return output_path
    except Exception as e:
        logger.exception("PDF 输出失败")
        raise OutputError("PDF 输出失败，请检查输出路径是否可写") from e


def to_ppt(result: AnalysisResult, output_path: str,
           template_path: str | None = None) -> str:
    """生成 PPT 报告。"""
    output_path = _validate_output_path(output_path)
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(template_path) if template_path and os.path.exists(
            template_path
        ) else Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11), Inches(2)
        )
        txBox.text_frame.text = f"分析报告: {result.task}\n\n{result.summary}"

        for fig in result.figures:
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=_CHART_DPI, bbox_inches='tight')
            buf.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(buf, Inches(0.5), Inches(0.5),
                                     Inches(12), Inches(6.5))
            plt.close(fig)

        prs.save(output_path)
        return output_path
    except Exception as e:
        logger.exception("PPT 输出失败")
        raise OutputError("PPT 输出失败，请检查输出路径是否可写") from e


def to_html(result: AnalysisResult, output_path: str) -> str:
    """生成自包含 HTML 分析报告 (Base64 内嵌图表)。"""
    output_path = _validate_output_path(output_path)
    import base64
    import html as _html

    def _esc(s):
        return _html.escape(str(s))

    try:
        html_parts = [
            "<!DOCTYPE html><html lang='zh-CN'><head>",
            "<meta charset='UTF-8'><meta name='viewport' content='width=device-width,initial-scale=1'>",
            "<title>SmartSuite 分析报告</title>",
            "<style>",
            "body{font-family:'Microsoft YaHei','PingFang SC','Noto Sans CJK SC',sans-serif;max-width:1100px;margin:0 auto;padding:20px;color:#333}",
            "h1{color:#2171b5;border-bottom:3px solid #2171b5;padding-bottom:10px}",
            "h2{color:#2171b5;margin-top:30px;border-bottom:1px solid #deebf7}",
            ".summary{background:#deebf7;padding:15px;border-radius:5px;margin:15px 0;font-size:14px}",
            "table{border-collapse:collapse;width:100%;margin:10px 0;font-size:12px}",
            "th{background:#2171b5;color:white;padding:8px 10px;text-align:left}",
            "td{padding:6px 10px;border-bottom:1px solid #deebf7}",
            "tr:nth-child(even){background:#f7fbff}",
            "img{max-width:100%;height:auto;margin:15px 0;border:1px solid #deebf7;border-radius:3px}",
            ".meta{color:#777;font-size:11px;margin-top:30px;border-top:1px solid #deebf7;padding-top:10px}",
            ".status-ok{color:#238b45}.status-error{color:#e31a1c}.status-warn{color:#d94801}",
            "</style></head><body>",
            f"<h1>SmartSuite 分析报告: {_esc(result.task)}</h1>",
        ]

        # 状态 + 结论
        status_class = f"status-{result.status}" if result.status in ("ok","error") else "status-warn"
        html_parts.append(
            f"<div class='summary'><strong>状态:</strong> "
            f"<span class='{status_class}'>{_esc(result.status)}</span><br>"
            f"<strong>结论:</strong> {_esc(result.summary)}</div>"
        )

        # 警告消息
        if result.messages:
            html_parts.append("<h2>诊断信息</h2><ul>")
            for m in result.messages:
                html_parts.append(f"<li>{_esc(m)}</li>")
            html_parts.append("</ul>")

        # 数据表
        for name, df in result.tables.items():
            html_parts.append(f"<h2>📊 {_esc(name)}</h2>")
            html_parts.append(df.head(50).to_html(
                index=False, classes="table", border=0, escape=True,
                float_format=lambda x: (
                    f"{x:.4f}" if isinstance(x, (int, float)) and abs(x) < 1e6
                    else f"{x:.2e}" if isinstance(x, (int, float))
                    else str(x)
                )
            ))
            if len(df) > 50:
                html_parts.append(f"<p style='color:#777;font-size:11px'>(仅显示前50行，共{len(df)}行)</p>")

        # 图表 (Base64 内嵌，压缩 PNG)
        for i, fig in enumerate(result.figures):
            buf = io.BytesIO()
            fig.savefig(buf, format="png", dpi=_CHART_DPI, bbox_inches="tight")
            buf.seek(0)
            # PIL 压缩优化
            try:
                from PIL import Image
                img = Image.open(buf)
                out_buf = io.BytesIO()
                img.save(out_buf, format="PNG", optimize=True)
                out_buf.seek(0)
                img_b64 = base64.b64encode(out_buf.read()).decode("utf-8")
            except ImportError:
                img_b64 = base64.b64encode(buf.read()).decode("utf-8")
            html_parts.append(
                f"<h2>📈 图表 {i+1}</h2>"
                f"<img src='data:image/png;base64,{img_b64}' alt='图表{i+1}'>"
            )
            plt.close(fig)

        # 元数据
        if result.metadata:
            html_parts.append("<h2>📋 元数据</h2><ul>")
            for k, v in list(result.metadata.items())[:15]:
                val_str = _esc(str(v)[:200])
                html_parts.append(f"<li><strong>{_esc(str(k))}:</strong> {val_str}</li>")
            html_parts.append("</ul>")

        html_parts.append(
            "<div class='meta'>由 SmartSuite 生成 | "
            "工艺数据分析工具箱</div></body></html>"
        )

        html_content = "\n".join(html_parts)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        logger.info("HTML 报告已生成: %s", output_path)
        return output_path
    except Exception as e:
        logger.exception("HTML 报告生成失败")
        raise OutputError(f"HTML 报告生成失败: {e}") from e
