"""Reporter — 多格式报告输出：Excel 图表 / PDF / PPT。"""
import io
import os

from smartexcel.core.contracts import AnalysisResult
from smartexcel.core.exceptions import OutputError


def to_excel(result: AnalysisResult, workbook,
             sheet_name: str = "分析结果") -> str:
    """将分析结果写入 Excel 新 Sheet。"""
    try:
        ws = workbook.sheets.add(sheet_name, after=workbook.sheets[-1])
        r = 1
        ws.range(f"A{r}").value = "分析结论"
        ws.range(f"A{r}").font.bold = True
        r += 1
        ws.range(f"A{r}").value = result.summary
        r += 2
        for name, df in result.tables.items():
            ws.range(f"A{r}").value = name
            ws.range(f"A{r}").font.bold = True
            r += 1
            ws.range(f"A{r}").value = df
            r += len(df) + 2
        for i, fig in enumerate(result.figures):
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
            buf.seek(0)
            pic = workbook.sheets.add(f"图表_{i + 1}", after=workbook.sheets[-1])
            pic.pictures.add(buf, left=pic.range("A1").left,
                             top=pic.range("A1").top, width=600, height=450)
        return sheet_name
    except Exception as e:
        raise OutputError(f"Excel 输出失败: {e}") from e


def to_pdf(result: AnalysisResult, output_path: str) -> str:
    """生成 PDF 报告。"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas

        c = canvas.Canvas(output_path, pagesize=A4)
        w, h = A4
        y = h - 50
        c.setFont("Helvetica-Bold", 16)
        c.drawString(50, y, f"分析报告: {result.task}")
        y -= 30
        c.setFont("Helvetica", 11)
        c.drawString(50, y, result.summary)
        y -= 50

        for name, df in list(result.tables.items())[:5]:
            if y < 150:
                c.showPage()
                y = h - 50
            c.setFont("Helvetica-Bold", 10)
            c.drawString(50, y, name)
            y -= 18
            c.setFont("Helvetica", 8)
            for _, row in df.head(15).iterrows():
                c.drawString(55, y, str(row.to_dict()))
                y -= 12

        for fig in result.figures[:3]:
            if y < 350:
                c.showPage()
                y = h - 50
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=100, bbox_inches='tight')
            buf.seek(0)
            c.drawImage(ImageReader(buf), 50, y - 300, width=450, height=300)
            y -= 320

        c.save()
        return output_path
    except Exception as e:
        raise OutputError(f"PDF 输出失败: {e}") from e


def to_ppt(result: AnalysisResult, output_path: str,
           template_path: str | None = None) -> str:
    """生成 PPT 报告。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(template_path) if template_path and os.path.exists(
            template_path
        ) else Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11), Inches(2)
        )
        txBox.text_frame.text = f"分析报告: {result.task}\n\n{result.summary}"

        for fig in result.figures:
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
            buf.seek(0)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(buf, Inches(0.5), Inches(0.5),
                                     Inches(12), Inches(6.5))

        prs.save(output_path)
        return output_path
    except Exception as e:
        raise OutputError(f"PPT 输出失败: {e}") from e
